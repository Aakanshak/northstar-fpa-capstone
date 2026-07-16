"""Build an editable eight-slide Q3 board deck from calculated workbook cells."""
from __future__ import annotations
import subprocess
from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT=Path(__file__).resolve().parents[1]; MODEL=ROOT/'03-excel-model'/'northstar_fpa_model.xlsx'
PPTX=Path(__file__).resolve().parent/'northstar_q3_board_deck.pptx'; PDF=Path(__file__).resolve().parent/'northstar_q3_board_deck.pdf'
NAVY=RGBColor(23,50,77); BLUE=RGBColor(47,117,181); TEAL=RGBColor(0,166,166); GREEN=RGBColor(112,173,71); RED=RGBColor(192,0,0); GOLD=RGBColor(244,177,131); LIGHT=RGBColor(239,244,248); GRAY=RGBColor(91,105,120); WHITE=RGBColor(255,255,255)

def money(v): return f"${v/1e6:.1f}M"
def add_title(slide,title,kicker=None):
    if kicker:
        t=slide.shapes.add_textbox(Inches(.65),Inches(.25),Inches(4),Inches(.28)).text_frame
        t.text=kicker.upper(); t.paragraphs[0].font.size=Pt(10); t.paragraphs[0].font.bold=True; t.paragraphs[0].font.color.rgb=TEAL
    box=slide.shapes.add_textbox(Inches(.65),Inches(.55),Inches(12),Inches(.62)); p=box.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(30); p.font.bold=True; p.font.color.rgb=NAVY
    line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.65),Inches(1.25),Inches(.65),Inches(.05)); line.fill.solid(); line.fill.fore_color.rgb=TEAL; line.line.fill.background()
def add_footer(slide,n):
    tb=slide.shapes.add_textbox(Inches(11.7),Inches(7.12),Inches(.9),Inches(.2)); p=tb.text_frame.paragraphs[0]; p.text=f"NORTHSTAR  |  {n}"; p.font.size=Pt(8); p.font.color.rgb=GRAY; p.alignment=PP_ALIGN.RIGHT
def textbox(slide,text,x,y,w,h,size=18,color=NAVY,bold=False):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=sh.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold; return sh
def metric(slide,label,value,x,y,color=BLUE):
    textbox(slide,value,x,y,2.8,.6,30,color,True); textbox(slide,label,x,y+.55,2.8,.4,12,GRAY,False)
def chart(slide,chart_type,categories,series,x,y,w,h,legend=True):
    data=CategoryChartData(); data.categories=categories
    for name,values in series: data.add_series(name,values)
    ch=slide.shapes.add_chart(chart_type,Inches(x),Inches(y),Inches(w),Inches(h),data).chart
    ch.has_legend=legend
    if legend: ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.font.size=Pt(10)
    ch.has_title=False; ch.value_axis.has_major_gridlines=True; ch.value_axis.tick_labels.font.size=Pt(9); ch.category_axis.tick_labels.font.size=Pt(9)
    colors=[BLUE,TEAL,GREEN,GOLD,RED]
    for i,s in enumerate(ch.series): s.format.fill.solid(); s.format.fill.fore_color.rgb=colors[i%len(colors)]
    return ch
def notes(slide,items,x=8.8,y=1.65,w=3.8):
    for i,(head,body) in enumerate(items):
        textbox(slide,head,x,y+i*1.45,w,.35,16,NAVY,True); textbox(slide,body,x,y+.38+i*1.45,w,.85,13,GRAY)

def build():
    wb=load_workbook(MODEL,data_only=True,read_only=True)
    rf=wb['Revenue Forecast']; arr=wb['ARR Bridge']; var=wb['Budget vs Actual']; met=wb['SaaS Metrics Dashboard']; bs=wb['Balance Sheet']; sens=wb['Sensitivity Analysis']
    months=[arr.cell(3,c).value.strftime('%b') for c in range(3,15)]
    base=[rf.cell(11,c).value for c in range(3,15)]; up=[rf.cell(20,c).value for c in range(3,15)]; down=[rf.cell(29,c).value for c in range(3,15)]
    ending_arr=rf['N34'].value; fy_rev=sum(rf.cell(33,c).value for c in range(3,15)); end_cash=bs['N5'].value; runway=met['N14'].value
    q3=sum(var.cell(8,c).value for c in range(9,12)); opex=sum(var.cell(18,c).value for c in range(3,15))
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    # 1 Executive summary
    s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
    textbox(s,'Q3 2026 BOARD UPDATE',.7,.55,5,.35,13,TEAL,True); textbox(s,'Growth remains intact.\nEfficiency now needs a gate.',.7,1.2,8.2,1.55,34,WHITE,True)
    metric(s,'FY2026 Base revenue',money(fy_rev),.75,3.25,TEAL); metric(s,'Ending ARR',money(ending_arr),3.75,3.25,GREEN); metric(s,'Ending runway',f'{runway:.1f} months',6.75,3.25,GOLD)
    textbox(s,f'Q3 revenue missed budget by {money(abs(q3))}; Sales and Marketing spend is ahead of plan. Hold incremental GTM hiring until conversion improves.',.75,5.15,10.7,1.1,18,WHITE)
    add_footer(s,1)
    # 2 ARR bridge
    s=prs.slides.add_slide(blank); add_title(s,'ARR reaches $58.3M, but growth depends on pipeline conversion','ARR BRIDGE')
    chart(s,XL_CHART_TYPE.COLUMN_CLUSTERED,months,[('New ARR',[arr.cell(5,c).value/1e6 for c in range(3,15)]),('Expansion',[arr.cell(6,c).value/1e6 for c in range(3,15)]),('Churn & contraction',[-(arr.cell(7,c).value+arr.cell(8,c).value)/1e6 for c in range(3,15)])],.65,1.55,7.8,4.9,True)
    notes(s,[('40.0M opening ARR','The source ledger is calibrated to the stated company scale.'),('Retention is stable',f'Dec NRR is {met["N6"].value:.1%}; expansion broadly offsets churn.'),('Pipeline is the swing factor','Base realizes 55% of stage-weighted pipeline after probability weighting.')]); add_footer(s,2)
    # 3 variance
    s=prs.slides.add_slide(blank); add_title(s,'Q3 miss is volume-led; GTM spend has not yet adjusted','BUDGET VS ACTUAL')
    chart(s,XL_CHART_TYPE.COLUMN_CLUSTERED,['Jul','Aug','Sep'],[('Volume',[var.cell(5,c).value/1e3 for c in range(9,12)]),('Price / mix',[var.cell(6,c).value/1e3 for c in range(9,12)])],.65,1.55,7.8,4.8,True)
    notes(s,[('$437K Q3 miss','July and August account for most of the shortfall.'),('$871K FY opex overrun','Marketing and Sales lead the unfavorable variance.'),('Interpretation','Pricing is resilient; close volume and timing remain the issue.')]); add_footer(s,3)
    # 4 metrics
    s=prs.slides.add_slide(blank); add_title(s,'Retention is healthy; efficiency is just below the Rule of 40','SAAS METRICS')
    metric(s,'December NRR',f'{met["N6"].value:.1%}',.8,1.75,GREEN); metric(s,'December GRR',f'{met["N7"].value:.1%}',3.8,1.75,GREEN); metric(s,'LTV : CAC',f'{met["N10"].value:.1f}x',6.8,1.75,TEAL); metric(s,'Rule of 40',f'{met["N11"].value:.1%}',9.8,1.75,GOLD)
    textbox(s,'The model points to a growth-efficiency tradeoff: retention is not the immediate constraint, but the company must improve new-logo conversion before adding GTM capacity.',.8,3.4,11.7,1.0,22,NAVY,True)
    textbox(s,'Management trigger',.8,5.0,2.5,.35,15,TEAL,True); textbox(s,'Release incremental hiring only after two consecutive months of improved conversion and shorter sales cycles.',.8,5.45,10.8,.8,18,GRAY)
    add_footer(s,4)
    # 5 scenarios
    s=prs.slides.add_slide(blank); add_title(s,'The FY revenue range is $48.2M-$55.9M','FY FORECAST')
    chart(s,XL_CHART_TYPE.LINE_MARKERS,months,[('Base',[v/1e6 for v in base]),('Upside',[v/1e6 for v in up]),('Downside',[v/1e6 for v in down])],.65,1.55,8.2,4.8,True)
    notes(s,[('Base | '+money(sum(base)),'1.1% churn, 0.55x pipeline factor, approved hiring pace.'),('Upside | '+money(sum(up)),'Lower churn, faster cycles, 0.70x pipeline realization.'),('Downside | '+money(sum(down)),'Higher churn, slower cycles, 0.38x pipeline realization.')],9.05,1.55,3.6); add_footer(s,5)
    # 6 cash sensitivity
    s=prs.slides.add_slide(blank); add_title(s,'Base runway is 12.3 months; hiring pace is the controllable lever','CASH & SENSITIVITY')
    cats=[sens.cell(r,2).value for r in range(25,29)]; lows=[sens.cell(r,3).value for r in range(25,29)]; highs=[sens.cell(r,5).value for r in range(25,29)]
    chart(s,XL_CHART_TYPE.BAR_CLUSTERED,cats,[('Low case',lows),('High case',highs)],.65,1.55,7.9,4.8,True)
    notes(s,[('Ending cash | '+money(end_cash),'Base case retains limited cushion.'),('Downside trigger','Move to 0.82x hiring if conversion and churn both miss.'),('Circularity control','Interest uses beginning cash; no iterative calculation is required.')]); add_footer(s,6)
    # 7 spend
    s=prs.slides.add_slide(blank); add_title(s,'Marketing and Sales account for the FY opex overrun','DEPARTMENT SPEND')
    depts=['Sales','Customer Success','Marketing','Product & Eng.','G&A']; vals=[]
    for r in range(12,17): vals.append(sum(var.cell(r,c).value for c in range(3,15))/1e3)
    chart(s,XL_CHART_TYPE.BAR_CLUSTERED,depts,[('Budget variance ($K)',vals)],.65,1.55,8.0,4.8,False)
    notes(s,[('Marketing | +$550K','Events and demand generation exceeded plan.'),('Sales | +$454K','Hiring and commissions ran ahead of conversion.'),('Offsets | -$402K','Customer Success and G&A discipline reduced the net overrun.')]); add_footer(s,7)
    # 8 actions
    s=prs.slides.add_slide(blank); add_title(s,'Protect the growth option without pre-funding unproven capacity','RISKS & RECOMMENDATIONS')
    actions=[('1','Gate GTM hiring','Require two months of conversion and cycle improvement before releasing roles.'),('2','Reallocate $300K','Shift low-attribution event spend toward partner and expansion programs.'),('3','Operationalize downside triggers','Review churn, conversion, DSO, and ending cash monthly; move to 0.82x hiring when two miss.')]
    for i,(num,head,body) in enumerate(actions):
        y=1.6+i*1.55; textbox(s,num,.8,y,.55,.55,26,TEAL,True); textbox(s,head,1.55,y,3.1,.45,19,NAVY,True); textbox(s,body,4.7,y,7.5,.85,17,GRAY)
    textbox(s,'BOARD ASK',.8,6.25,1.6,.3,12,TEAL,True); textbox(s,'Approve the Base case with a conditional GTM hiring gate and monthly downside trigger review.',2.25,6.1,9.8,.6,20,NAVY,True); add_footer(s,8)
    for slide in prs.slides:
        if slide.background.fill.type is None: slide.background.fill.solid(); slide.background.fill.fore_color.rgb=WHITE
    prs.save(PPTX); export_pdf(PPTX,PDF); print(f'Built {PPTX}\nBuilt {PDF}')

def export_pdf(pptx,pdf):
    ps=f'''$ppt=New-Object -ComObject PowerPoint.Application; $pres=$ppt.Presentations.Open("{pptx}",$false,$false,$false); $pres.SaveAs("{pdf}",32); $pres.Close(); $ppt.Quit(); [Runtime.InteropServices.Marshal]::ReleaseComObject($pres)|Out-Null; [Runtime.InteropServices.Marshal]::ReleaseComObject($ppt)|Out-Null'''
    subprocess.run(['powershell','-NoProfile','-Command',ps],check=True)

if __name__=='__main__': build()
